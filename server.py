from flask import Flask, request, send_from_directory, Response
from flask_cors import CORS
import os, subprocess, tempfile, traceback, threading

app = Flask(__name__)
CORS(app)

@app.route('/')
def index():
    return send_from_directory('.', 'index.html')

@app.route('/tool.html')
def tool_page():
    return send_from_directory('.', 'tool.html')

def cleanup_later(tmp_dir, delay=60):
    def run():
        import time
        time.sleep(delay)
        import shutil
        try: shutil.rmtree(tmp_dir, ignore_errors=True)
        except: pass
    threading.Thread(target=run, daemon=True).start()

def serve_file(filepath, download_name):
    with open(filepath, 'rb') as f:
        data = f.read()
    return Response(
        data,
        headers={
            'Content-Type': 'application/octet-stream',
            'Content-Disposition': f'attachment; filename="{download_name}"',
            'Content-Length': str(len(data)),
            'Access-Control-Expose-Headers': 'Content-Disposition, Content-Length'
        }
    )

def libreoffice_convert(input_path, output_dir, output_ext):
    cmd = [
        'libreoffice', '--headless', '--norestore',
        '--convert-to', output_ext,
        '--outdir', output_dir,
        input_path
    ]
    result = subprocess.run(cmd, capture_output=True, timeout=120)
    if result.returncode != 0:
        raise RuntimeError(f'LibreOffice error: {result.stderr.decode("utf-8", errors="replace")}')
    base = os.path.splitext(os.path.basename(input_path))[0]
    output_path = os.path.join(output_dir, base + '.' + output_ext)
    if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
        raise RuntimeError('Conversion produced empty file')
    return output_path

@app.route('/convert/pptx', methods=['POST'])
def api_convert_pptx():
    if 'file' not in request.files:
        return 'No file uploaded', 400
    file = request.files['file']
    if not file.filename.lower().endswith('.pptx'):
        return 'Not a PPTX file', 400
    tmp_dir = tempfile.mkdtemp()
    input_path = os.path.join(tmp_dir, file.filename)
    try:
        file.save(input_path)
        output_path = libreoffice_convert(input_path, tmp_dir, 'pdf')
        download_name = file.filename.rsplit('.', 1)[0] + '.pdf'
        resp = serve_file(output_path, download_name)
        cleanup_later(tmp_dir)
        return resp
    except Exception as e:
        traceback.print_exc()
        cleanup_later(tmp_dir)
        return str(e), 500

@app.route('/convert/xlsx', methods=['POST'])
def api_convert_xlsx():
    if 'file' not in request.files:
        return 'No file uploaded', 400
    file = request.files['file']
    fname = file.filename.lower()
    if not (fname.endswith('.xlsx') or fname.endswith('.xls')):
        return 'Not an Excel file', 400
    tmp_dir = tempfile.mkdtemp()
    input_path = os.path.join(tmp_dir, file.filename)
    try:
        file.save(input_path)
        output_path = libreoffice_convert(input_path, tmp_dir, 'pdf')
        download_name = file.filename.rsplit('.', 1)[0] + '.pdf'
        resp = serve_file(output_path, download_name)
        cleanup_later(tmp_dir)
        return resp
    except Exception as e:
        traceback.print_exc()
        cleanup_later(tmp_dir)
        return str(e), 500

@app.route('/convert/pdf-to-docx', methods=['POST'])
def api_convert_pdf_to_docx():
    if 'file' not in request.files:
        return 'No file uploaded', 400
    file = request.files['file']
    if not file.filename.lower().endswith('.pdf'):
        return 'Not a PDF file', 400
    tmp_dir = tempfile.mkdtemp()
    input_path = os.path.join(tmp_dir, file.filename)
    output_path = os.path.join(tmp_dir, file.filename.rsplit('.', 1)[0] + '.docx')
    try:
        file.save(input_path)
        from pdf2docx import Converter
        cv = Converter(input_path)
        cv.convert(output_path)
        cv.close()
        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            return 'Conversion produced empty file', 500
        download_name = file.filename.rsplit('.', 1)[0] + '.docx'
        resp = serve_file(output_path, download_name)
        cleanup_later(tmp_dir)
        return resp
    except Exception as e:
        traceback.print_exc()
        cleanup_later(tmp_dir)
        return str(e), 500

@app.route('/convert/pdf-to-pptx', methods=['POST'])
def api_convert_pdf_to_pptx():
    if 'file' not in request.files:
        return 'No file uploaded', 400
    file = request.files['file']
    if not file.filename.lower().endswith('.pdf'):
        return 'Not a PDF file', 400
    tmp_dir = tempfile.mkdtemp()
    input_path = os.path.join(tmp_dir, file.filename)
    output_path = os.path.join(tmp_dir, file.filename.rsplit('.', 1)[0] + '.pptx')
    try:
        file.save(input_path)
        import fitz
        from pptx import Presentation
        from pptx.util import Inches
        pdf_doc = fitz.open(input_path)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            pix = page.get_pixmap(dpi=200)
            img_path = os.path.join(tmp_dir, f'page_{page_num}.png')
            pix.save(img_path)
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        pdf_doc.close()
        prs.save(output_path)
        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            return 'Conversion produced empty file', 500
        download_name = file.filename.rsplit('.', 1)[0] + '.pptx'
        resp = serve_file(output_path, download_name)
        cleanup_later(tmp_dir)
        return resp
    except Exception as e:
        traceback.print_exc()
        cleanup_later(tmp_dir)
        return str(e), 500

@app.route('/convert/protect', methods=['POST'])
def api_protect_pdf():
    if 'file' not in request.files:
        return 'No file uploaded', 400
    file = request.files['file']
    password = request.form.get('password', '')
    if not file.filename.lower().endswith('.pdf'):
        return 'Not a PDF file', 400
    if not password or len(password) < 3:
        return 'Password must be at least 3 characters', 400
    tmp_dir = tempfile.mkdtemp()
    input_path = os.path.join(tmp_dir, file.filename)
    output_path = os.path.join(tmp_dir, file.filename.rsplit('.', 1)[0] + '_protected.pdf')
    try:
        file.save(input_path)
        import pikepdf
        pdf = pikepdf.open(input_path)
        pdf.save(
            output_path,
            encryption=pikepdf.Encryption(
                owner=password,
                user=password,
                R=6,
            )
        )
        pdf.close()
        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            return 'Protection produced empty file', 500
        download_name = file.filename.rsplit('.', 1)[0] + '_protected.pdf'
        resp = serve_file(output_path, download_name)
        cleanup_later(tmp_dir)
        return resp
    except Exception as e:
        traceback.print_exc()
        cleanup_later(tmp_dir)
        return str(e), 500

@app.route('/convert/pdf-to-xlsx', methods=['POST'])
def api_convert_pdf_to_xlsx():
    if 'file' not in request.files:
        return 'No file uploaded', 400
    file = request.files['file']
    if not file.filename.lower().endswith('.pdf'):
        return 'Not a PDF file', 400
    tmp_dir = tempfile.mkdtemp()
    input_path = os.path.join(tmp_dir, file.filename)
    output_path = os.path.join(tmp_dir, file.filename.rsplit('.', 1)[0] + '.xlsx')
    try:
        file.save(input_path)
        import pdfplumber
        import fitz
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
        wb = Workbook()
        first_sheet = True
        pdf_doc = fitz.open(input_path)
        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            if first_sheet:
                ws = wb.active
                ws.title = f'Page {page_num + 1}'
                first_sheet = False
            else:
                ws = wb.create_sheet(title=f'Page {page_num + 1}')
            pdf_page = pdfplumber.open(input_path).pages[page_num]
            tables = pdf_page.extract_tables()
            if tables:
                for table in tables:
                    for row_idx, row in enumerate(table):
                        for col_idx, cell in enumerate(row):
                            val = str(cell).strip() if cell else ''
                            c = ws.cell(row=row_idx + 1, column=col_idx + 1, value=val)
                            c.alignment = Alignment(wrap_text=True, vertical='top')
                            c.border = Border(
                                left=Side(style='thin'), right=Side(style='thin'),
                                top=Side(style='thin'), bottom=Side(style='thin')
                            )
                            if row_idx == 0:
                                c.font = Font(bold=True)
                                c.fill = PatternFill(start_color='D9E1F2', end_color='D9E1F2', fill_type='solid')
                    for col in ws.columns:
                        max_len = max(len(str(cell.value or '')) for cell in col)
                        ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 40)
            else:
                words = page.get_text('words')
                if not words:
                    continue
                lines = {}
                for w in words:
                    y = round(w[1], 0)
                    found = False
                    for key in lines:
                        if abs(key - y) < 5:
                            lines[key].append(w)
                            found = True
                            break
                    if not found:
                        lines[y] = [w]
                all_x = sorted(set(round(w[0], 0) for w in words))
                gaps = [(all_x[i+1] - all_x[i], i) for i in range(len(all_x)-1)]
                gaps.sort(reverse=True)
                col_positions = [all_x[0]]
                for g, i in gaps:
                    if g > 50:
                        col_positions.append(all_x[i + 1])
                col_positions.sort()
                for ci in range(len(col_positions)):
                    ws.column_dimensions[chr(65 + ci)].width = 20
                sorted_ys = sorted(lines.keys(), reverse=True)
                for ri, y_key in enumerate(sorted_ys):
                    row_words = sorted(lines[y_key], key=lambda w: w[0])
                    row_data = {}
                    for w in row_words:
                        cx = round(w[0], 0)
                        best_col = 0
                        best_dist = abs(cx - col_positions[0])
                        for ci, cp in enumerate(col_positions):
                            dist = abs(cx - cp)
                            if dist < best_dist:
                                best_dist = dist
                                best_col = ci
                        if best_col not in row_data:
                            row_data[best_col] = []
                        row_data[best_col].append(w[4])
                    for col_idx in sorted(row_data.keys()):
                        val = ' '.join(row_data[col_idx])
                        cell = ws.cell(row=ri + 1, column=col_idx + 1, value=val)
                        cell.alignment = Alignment(wrap_text=True, vertical='top')
                        cell.border = Border(
                            left=Side(style='hair'), right=Side(style='hair'),
                            top=Side(style='hair'), bottom=Side(style='hair')
                        )
            pdf_page.close()
        pdf_doc.close()
        wb.save(output_path)
        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            return 'Conversion produced empty file', 500
        download_name = file.filename.rsplit('.', 1)[0] + '.xlsx'
        resp = serve_file(output_path, download_name)
        cleanup_later(tmp_dir)
        return resp
    except Exception as e:
        traceback.print_exc()
        cleanup_later(tmp_dir)
        return str(e), 500

@app.route('/<path:file>')
def static_files(file):
    return send_from_directory('.', file)

if __name__ == '__main__':
    print('Server running at http://localhost:8080')
    app.run(host='0.0.0.0', port=8080, debug=False)
